import tempfile
import unittest
from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Pt

from utils.ppt_tables_slide_24 import (
    SLIDE24_SHEET_NAME,
    SLIDE24_TABLE_ALT_TEXT,
    apply_slide24_table_file,
    apply_slide24_table_to_presentation,
    extract_slide24_table_headers,
    extract_slide24_table_values,
)


class TestPptTablesSlide24(unittest.TestCase):
    def _create_workbook(self, xlsx_path: Path) -> None:
        wb = Workbook()
        ws = wb.active
        ws.title = SLIDE24_SHEET_NAME

        ws["C3"] = "IGNORAR NO PPT"
        ws["D3"] = "4T24"
        ws["E3"] = "3T25"
        ws["F3"] = "4T25"
        ws["G3"] = 2024
        ws["H3"] = 2025
        ws["I3"] = "Variação"
        ws["I4"] = "4T25/3T25"
        ws["J4"] = "4T25/4T24"
        ws["K4"] = "2025/2024"

        ws["C5"] = "Receitas Totais"
        ws["D5"] = 3214
        ws["E5"] = 2919
        ws["F5"] = 3136
        ws["G5"] = 11980
        ws["H5"] = 11913
        ws["I5"] = 0.074
        ws["I5"].number_format = "0.0%"
        ws["J5"] = -0.024
        ws["J5"].number_format = "0.0%"
        ws["K5"] = -0.006
        ws["K5"].number_format = "0.0%"

        ws["C10"] = "Custo de crédito"
        ws["D10"] = -776
        ws["E10"] = -880
        ws["F10"] = -1029
        ws["G10"] = -3593
        ws["H10"] = -3698
        ws["I10"] = 16.8
        ws["J10"] = 32.6
        ws["K10"] = 2.9

        wb.save(xlsx_path)

    def _set_alt_text(self, shape, alt_text: str) -> None:
        cnv = shape._element.xpath(".//p:cNvPr")
        self.assertTrue(cnv)
        cnv[0].set("descr", alt_text)

    def _create_presentation_with_table(self, pptx_path: Path) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_table(16, 9, 0, 0, 9_000_000, 4_500_000)
        self._set_alt_text(shape, SLIDE24_TABLE_ALT_TEXT)
        table = shape.table

        table.cell(0, 0).merge(table.cell(1, 0))
        for col_idx in range(1, 6):
            table.cell(0, col_idx).merge(table.cell(1, col_idx))
        table.cell(0, 6).merge(table.cell(0, 8))

        for row_idx in range(16):
            for col_idx in range(9):
                cell = table.cell(row_idx, col_idx)
                if getattr(cell, "is_spanned", False):
                    continue
                cell.text = f"placeholder-{row_idx}-{col_idx}"

        table.cell(0, 0).text = "FIXED PPT"
        styled_run = table.cell(2, 1).text_frame.paragraphs[0].runs[0]
        styled_run.font.size = Pt(22)
        styled_run.font.bold = True
        prs.save(pptx_path)

    def test_extract_slide24_table_ranges_format_expected_columns(self):
        with tempfile.TemporaryDirectory() as td:
            xlsx_path = Path(td) / "test.xlsx"
            self._create_workbook(xlsx_path)

            headers = extract_slide24_table_headers(xlsx_path=xlsx_path)
            values = extract_slide24_table_values(xlsx_path=xlsx_path)

        self.assertEqual(headers[0][0], "IGNORAR NO PPT")
        self.assertEqual(headers[0][1], "4T24")
        self.assertEqual(headers[0][6], "Variação")
        self.assertEqual(headers[1][6], "4T25/3T25")
        self.assertEqual(values[0][0], "Receitas Totais")
        self.assertEqual(values[0][1], "3.214")
        self.assertEqual(values[0][6], "7,4")
        self.assertEqual(values[0][7], "-2,4")
        self.assertEqual(values[0][8], "-0,6")
        self.assertEqual(values[5][1], "(776)")

    def test_apply_slide24_table_file_writes_values_and_preserves_style(self):
        with tempfile.TemporaryDirectory() as td:
            tmpdir = Path(td)
            xlsx_path = tmpdir / "test.xlsx"
            pptx_input = tmpdir / "input.pptx"
            pptx_output = tmpdir / "output.pptx"
            self._create_workbook(xlsx_path)
            self._create_presentation_with_table(pptx_input)

            result = apply_slide24_table_file(
                pptx_path=pptx_input,
                output_path=pptx_output,
                xlsx_path=xlsx_path,
            )

            prs = Presentation(str(pptx_output))
            table_shape = next(shape for shape in prs.slides[0].shapes if getattr(shape, "has_table", False))
            table = table_shape.table

        self.assertTrue(result.found)
        self.assertEqual(result.slide_index, 1)
        self.assertEqual(result.skipped_fixed_cells, 1)
        self.assertEqual(result.skipped_spanned_cells, 8)
        self.assertEqual(table.cell(0, 0).text, "FIXED PPT")
        self.assertEqual(table.cell(0, 1).text, "4T24")
        self.assertEqual(table.cell(0, 6).text, "Variação")
        self.assertEqual(table.cell(1, 6).text, "4T25/3T25")
        self.assertEqual(table.cell(1, 7).text, "4T25/4T24")
        self.assertEqual(table.cell(1, 8).text, "2025/2024")
        self.assertEqual(table.cell(2, 0).text, "Receitas Totais")
        self.assertEqual(table.cell(2, 1).text, "3.214")
        self.assertEqual(table.cell(2, 6).text, "7,4")
        self.assertEqual(table.cell(7, 1).text, "(776)")
        styled_run = table.cell(2, 1).text_frame.paragraphs[0].runs[0]
        self.assertEqual(int(round(styled_run.font.size.pt)), 22)
        self.assertTrue(styled_run.font.bold)

    def test_apply_slide24_table_to_presentation_returns_not_found_without_strict(self):
        with tempfile.TemporaryDirectory() as td:
            xlsx_path = Path(td) / "test.xlsx"
            self._create_workbook(xlsx_path)
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])

            result = apply_slide24_table_to_presentation(prs, xlsx_path=xlsx_path)

        self.assertFalse(result.found)
        self.assertEqual(result.written_cells, 0)


if __name__ == "__main__":
    unittest.main()
