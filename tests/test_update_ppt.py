import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Pt

from update_ppt import _replace_text_in_shape, update_presentation
from utils.ppt_tables_slide_24 import SLIDE24_TABLE_ALT_TEXT


class TestUpdatePpt(unittest.TestCase):
    def _make_text_shape(self, text: str):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
        shape.text_frame.text = text
        return shape

    def test_replace_text_in_shape_formats_var_pp_fields_as_pp(self):
        shape = self._make_text_shape("{{VAR_TEST}}")

        replaced = _replace_text_in_shape(
            shape,
            {"VAR_TEST": "-0.9"},
            pp_field_ids={"VAR_TEST"},
        )

        joined = "".join(run.text for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)
        self.assertEqual(replaced, 1)
        self.assertEqual(joined, "▼ 0,9 p.p.")

    def test_replace_text_in_shape_keeps_percent_logic_for_regular_var(self):
        shape = self._make_text_shape("{{VAR_TEST}}")

        replaced = _replace_text_in_shape(
            shape,
            {"VAR_TEST": "-0.9"},
            pp_field_ids=set(),
        )

        joined = "".join(run.text for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)
        self.assertEqual(replaced, 1)
        self.assertEqual(joined, "▼ 90,0%")

    def test_update_presentation_applies_slide24_table_when_xlsx_is_provided(self):
        with tempfile.TemporaryDirectory() as td:
            tmpdir = Path(td)
            pptx_path = tmpdir / "input.pptx"
            output_path = tmpdir / "output.pptx"
            xlsx_path = tmpdir / "input.xlsx"

            wb = Workbook()
            ws = wb.active
            ws.title = "DRE Saida"
            ws["C3"] = "FIXO EXCEL"
            ws["D3"] = "4T24"
            ws["E3"] = "3T25"
            ws["F3"] = "4T25"
            ws["G3"] = 2024
            ws["H3"] = 2025
            ws["I3"] = "Variação"
            ws["I4"] = "4T25/3T25"
            ws["J4"] = "4T25/4T24"
            ws["K4"] = "2025/2024"
            ws["C5"] = "Receitas Totais"
            ws["D5"] = 3214
            ws["I5"] = 0.074
            ws["I5"].number_format = "0.0%"
            wb.save(xlsx_path)

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_table(16, 9, 0, 0, 9_000_000, 4_500_000)
            cnv = shape._element.xpath(".//p:cNvPr")
            self.assertTrue(cnv)
            cnv[0].set("descr", SLIDE24_TABLE_ALT_TEXT)
            table = shape.table
            table.cell(0, 0).merge(table.cell(1, 0))
            for col_idx in range(1, 6):
                table.cell(0, col_idx).merge(table.cell(1, col_idx))
            table.cell(0, 6).merge(table.cell(0, 8))
            for r in range(16):
                for c in range(9):
                    cell = table.cell(r, c)
                    if getattr(cell, "is_spanned", False):
                        continue
                    cell.text = f"placeholder-{r}-{c}"
            table.cell(0, 0).text = "FIXED PPT"
            run = table.cell(2, 1).text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(20)
            run.font.bold = True
            prs.save(pptx_path)

            update_presentation(
                pptx_path=pptx_path,
                output_path=output_path,
                images_dir=tmpdir,
                allow_placeholder_text=False,
                text_json=None,
                xlsx_path=xlsx_path,
            )

            updated = Presentation(str(output_path))
            updated_table = next(shape for shape in updated.slides[0].shapes if getattr(shape, "has_table", False)).table

        self.assertEqual(updated_table.cell(0, 0).text, "FIXED PPT")
        self.assertEqual(updated_table.cell(0, 1).text, "4T24")
        self.assertEqual(updated_table.cell(0, 6).text, "Variação")
        self.assertEqual(updated_table.cell(1, 6).text, "4T25/3T25")
        self.assertEqual(updated_table.cell(2, 0).text, "Receitas Totais")
        self.assertEqual(updated_table.cell(2, 1).text, "3.214")
        self.assertEqual(updated_table.cell(2, 6).text, "7,4")
        styled_run = updated_table.cell(2, 1).text_frame.paragraphs[0].runs[0]
        self.assertEqual(int(round(styled_run.font.size.pt)), 20)
        self.assertTrue(styled_run.font.bold)

    def test_update_presentation_keeps_updating_when_slide24_table_fails(self):
        with tempfile.TemporaryDirectory() as td:
            tmpdir = Path(td)
            pptx_path = tmpdir / "input.pptx"
            output_path = tmpdir / "output.pptx"
            xlsx_path = tmpdir / "input.xlsx"

            wb = Workbook()
            wb.active.title = "DRE Saida"
            wb.save(xlsx_path)

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(pptx_path)

            with patch("update_ppt.apply_slide24_table_to_presentation", side_effect=ValueError("layout incompatível")):
                update_presentation(
                    pptx_path=pptx_path,
                    output_path=output_path,
                    images_dir=tmpdir,
                    allow_placeholder_text=False,
                    text_json=None,
                    xlsx_path=xlsx_path,
                )

            self.assertTrue(output_path.exists())


if __name__ == "__main__":
    unittest.main()
