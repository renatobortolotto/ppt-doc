from __future__ import annotations

import argparse
import json
import logging
import re
import warnings
from pathlib import Path

from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _replace_shape_with_picture(slide, shape, image_path: Path) -> None:
    left, top, box_w, box_h = shape.left, shape.top, shape.width, shape.height

    # Preserve image aspect ratio: fit inside the box and center.
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    if img_w <= 0 or img_h <= 0:
        new_left, new_top, new_w, new_h = left, top, box_w, box_h
    else:
        scale = min(float(box_w) / float(img_w), float(box_h) / float(img_h))
        new_w = int(round(img_w * scale))
        new_h = int(round(img_h * scale))
        new_left = int(round(left + (box_w - new_w) / 2))
        new_top = int(round(top + (box_h - new_h) / 2))

    slide.shapes.add_picture(str(image_path), new_left, new_top, width=new_w, height=new_h)

    # Remove old shape
    el = shape._element
    el.getparent().remove(el)


def _get_shape_alt_text(shape) -> str | None:
    try:
        cnv = shape._element.xpath('.//p:cNvPr')
        if cnv:
            return cnv[0].get('descr')
    except Exception:
        return None
    return None


def _replace_picture_image_in_place(slide, picture_shape, image_path: Path) -> None:
    """Replace a PICTURE's image without changing its geometry/crop.

    Uses python-pptx internals to swap the blip rId to a new image part.
    """
    if not image_path.exists():
        raise FileNotFoundError(f"Imagem não encontrada: {image_path}")

    # Add or reuse image part, then rel to slide, returning rId.
    image_part, rId = slide.part.get_or_add_image_part(str(image_path))

    # Swap embed reference (keeps extents/crop/position from existing picture).
    blips = picture_shape._element.xpath('.//a:blip')
    if not blips:
        raise ValueError(f"Não achei <a:blip> para substituir na shape {picture_shape.name!r}")
    blips[0].set(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed',
        rId,
    )


def _flatten_text_payload(payload: object) -> dict[str, str]:
    """Extract a flat key->string map from an LLM JSON payload.

    Expected model output shape (example):
        {
          "titles": {"slide1_title": "..."},
          "subtitles": {"slide1_subtitle": "..."}
        }

    We also accept top-level keys directly.
    """

    def _camel_to_snake(key: str) -> str:
        # slide1Title -> slide1_title ; Slide1Title -> slide1_title
        s1 = re.sub(r"(.)([A-Z][a-z]+)", r"\1_\2", key)
        s2 = re.sub(r"([a-z0-9])([A-Z])", r"\1_\2", s1)
        return s2.replace("__", "_").lower()

    def _add(mapping: dict[str, str], key: str, value: object) -> None:
        if value is None:
            return
        text = str(value)
        mapping[key] = text

        snake = _camel_to_snake(key)
        # Only add the snake_case alias if it differs and doesn't already exist.
        if snake != key and snake not in mapping:
            mapping[snake] = text

    mapping: dict[str, str] = {}
    if not isinstance(payload, dict):
        return mapping

    for section_key in ("titles", "subtitles"):
        section = payload.get(section_key)
        if isinstance(section, dict):
            for k, v in section.items():
                _add(mapping, str(k), v)

    # Also accept any top-level string keys (optional)
    for k, v in payload.items():
        if k in ("titles", "subtitles"):
            continue
        if isinstance(v, str):
            _add(mapping, str(k), v)

    return mapping


def _replace_text_in_shape(shape, mapping: dict[str, str]) -> int:
    """Replace text placeholders inside a shape. Returns number of replacements."""

    def _is_var_field(key: str) -> bool:
        return bool(key) and str(key).upper().startswith("VAR_")

    def _parse_float_loose(text: str) -> float | None:
        if text is None:
            return None
        s = str(text).strip()
        if not s:
            return None
        # Common cases: "-1,2%", "+0.5%", "0.03".
        s = s.replace("%", "")
        s = s.replace(" ", "")
        # Keep only characters relevant for a float.
        s = re.sub(r"[^0-9eE+\-\.,]", "", s)
        # Prefer comma as decimal separator.
        if s.count(",") == 1 and s.count(".") == 0:
            s = s.replace(",", ".")
        # If both separators appear, assume dot is thousands and comma is decimal.
        if "," in s and "." in s:
            s = s.replace(".", "")
            s = s.replace(",", ".")
        try:
            return float(s)
        except Exception:
            return None

    def _format_var_indicator_parts(raw: str) -> tuple[str, str, RGBColor] | None:
        val = _parse_float_loose(raw)
        if val is None:
            return None

        # Decide if the source is already a percentage.
        raw_s = str(raw).strip()
        has_percent = "%" in raw_s

        eps = 1e-9
        if abs(val) <= eps:
            glyph = "●"
            color = RGBColor(0x7F, 0x7F, 0x7F)  # gray
            mag = 0.0
        elif val > 0:
            glyph = "▲"
            color = RGBColor(0x00, 0xB0, 0x50)  # green
            mag = abs(val)
        else:
            glyph = "▼"
            color = RGBColor(0xC0, 0x00, 0x00)  # red
            mag = abs(val)

        if has_percent:
            # Preserve user's/Excel's formatted magnitude as much as possible.
            # Just remove sign characters.
            cleaned = raw_s
            cleaned = cleaned.replace("+", "")
            cleaned = cleaned.replace("-", "")
            cleaned = cleaned.strip()
            return glyph, cleaned, color

        # If the cell is a fraction (0.0123), treat as percent.
        if abs(mag) <= 1.0:
            mag_pct = mag * 100.0
        else:
            mag_pct = mag

        # 1 decimal is a good default for quarter deltas.
        mag_txt = f"{mag_pct:.1f}%".replace(".", ",")
        return glyph, mag_txt, color

    def _set_shape_text_with_var(*, raw: str) -> bool:
        """Set shape text to 'glyph value' with only glyph colored.

        Returns True if it applied VAR formatting, False otherwise.
        """

        formatted = _format_var_indicator_parts(raw)
        if formatted is None:
            return False
        glyph, mag_txt, rgb = formatted

        # Capture the current styling so we don't lose template formatting.
        style_src = None
        try:
            p0 = shape.text_frame.paragraphs[0]
            if p0.runs:
                style_src = p0.runs[0]
        except Exception:
            style_src = None

        def _copy_font(src_run, dst_run) -> None:
            if src_run is None:
                return
            try:
                sf = src_run.font
                df = dst_run.font
                df.name = sf.name
                df.size = sf.size
                df.bold = sf.bold
                df.italic = sf.italic
                df.underline = sf.underline
                # Keep the existing color for the numeric run; glyph will override.
                try:
                    df.color.rgb = sf.color.rgb
                except Exception:
                    pass
            except Exception:
                return

        shape.text_frame.text = ""
        p = shape.text_frame.paragraphs[0]
        p.text = ""

        r1 = p.add_run()
        r1.text = glyph
        _copy_font(style_src, r1)
        try:
            r1.font.color.rgb = rgb
        except Exception:
            pass

        r2 = p.add_run()
        r2.text = f" {mag_txt}"
        _copy_font(style_src, r2)
        return True

    def _rebuild_paragraph_with_tokens(paragraph) -> int:
        """Rebuild paragraph text replacing tokens; VAR_* tokens get colored glyph only."""

        try:
            src = paragraph.text or ""
        except Exception:
            src = ""
        if not src:
            return 0

        # Only rebuild if the paragraph contains at least one VAR_* token.
        # For normal tokens, we keep the original run structure to preserve formatting.
        contains_var_token = False
        for m in re.finditer(r"\{\{([^}]+)\}\}", src):
            key = m.group(1)
            if key in mapping and _is_var_field(key):
                contains_var_token = True
                break
        if not contains_var_token:
            return 0

        # Capture styling from the first run (if any) so we don't lose template sizing.
        style_src = None
        try:
            if paragraph.runs:
                style_src = paragraph.runs[0]
        except Exception:
            style_src = None

        def _copy_font(src_run, dst_run) -> None:
            if src_run is None:
                return
            try:
                sf = src_run.font
                df = dst_run.font
                df.name = sf.name
                df.size = sf.size
                df.bold = sf.bold
                df.italic = sf.italic
                df.underline = sf.underline
                try:
                    df.color.rgb = sf.color.rgb
                except Exception:
                    pass
            except Exception:
                return

        # Build segments: (text, optional_rgb)
        segs: list[tuple[str, RGBColor | None]] = []
        replaced_local = 0
        pos = 0
        for m in re.finditer(r"\{\{([^}]+)\}\}", src):
            key = m.group(1)
            if key not in mapping:
                continue

            start, end = m.span(0)
            if start > pos:
                segs.append((src[pos:start], None))

            raw = mapping[key]
            if _is_var_field(key):
                parts = _format_var_indicator_parts(raw)
                if parts is None:
                    segs.append((raw, None))
                else:
                    glyph, mag_txt, rgb = parts
                    segs.append((glyph, rgb))
                    segs.append((f" {mag_txt}", None))
            else:
                segs.append((raw, None))

            pos = end
            replaced_local += 1

        if replaced_local == 0:
            return 0

        if pos < len(src):
            segs.append((src[pos:], None))

        paragraph.text = ""
        for text, rgb in segs:
            run = paragraph.add_run()
            run.text = text
            _copy_font(style_src, run)
            if rgb is not None:
                try:
                    run.font.color.rgb = rgb
                except Exception:
                    pass

        return replaced_local

    if not getattr(shape, "has_text_frame", False):
        return 0

    # Strategy:
    # 1) If Alt Text matches a key, replace whole text.
    # 2) Replace tokens {{key}} within runs (preserve formatting when possible).
    # 3) If the whole text equals a key, replace it.

    replaced = 0
    alt = _get_shape_alt_text(shape)
    if alt and alt in mapping:
        raw = mapping[alt]
        if _is_var_field(alt):
            if _set_shape_text_with_var(raw=raw):
                return 1
        shape.text_frame.text = raw
        return 1

    full_text = (shape.text_frame.text or "")
    if not full_text.strip():
        return 0

    # Whole-text placeholder (e.g., slide1_title)
    key = full_text.strip()
    if key in mapping:
        raw = mapping[key]
        if _is_var_field(key):
            if _set_shape_text_with_var(raw=raw):
                return 1
        shape.text_frame.text = raw
        return 1

    # Token replacement (recommended): {{slide1_title}}
    for paragraph in shape.text_frame.paragraphs:
        # Only rebuild when needed for VAR_* (mixed-color) formatting.
        replaced_here = _rebuild_paragraph_with_tokens(paragraph)
        if replaced_here:
            replaced += replaced_here
            continue

        # Try run-based replacement first (keeps formatting).
        for run in paragraph.runs:
            t = run.text
            if not t:
                continue
            for k, v in mapping.items():
                token = "{{" + k + "}}"
                if token in t:
                    run.text = t.replace(token, v)
                    t = run.text
                    replaced += 1

        # Fallback if token spans multiple runs
        try:
            p_text = paragraph.text
        except Exception:
            p_text = ""
        if p_text:
            new_text = p_text
            any_token = False
            for k, v in mapping.items():
                token = "{{" + k + "}}"
                if token in new_text:
                    new_text = new_text.replace(token, v)
                    any_token = True
            if any_token and new_text != p_text:
                paragraph.text = new_text
                replaced += 1

    return replaced


def update_presentation(
    pptx_path: Path,
    output_path: Path,
    images_dir: Path,
    allow_placeholder_text: bool,
    text_json: Path | None,
    text_payload: object | None = None,
) -> tuple[int, int, int, list[str], list[str], list[str]]:
    prs = Presentation(str(pptx_path))

    text_mapping: dict[str, str] = {}
    payload: object | None = None
    if text_payload is not None:
        payload = text_payload
    elif text_json is not None:
        payload = json.loads(text_json.read_text(encoding="utf-8"))

    if payload is not None:
        # Support either raw model output or wrapper: {"response": {...}}
        if isinstance(payload, dict) and "response" in payload and isinstance(payload["response"], dict):
            payload = payload["response"]
        text_mapping = _flatten_text_payload(payload)

    replaced_pictures = 0
    replaced_placeholders = 0
    replaced_text = 0
    missing_files: list[str] = []
    replaced_files: list[str] = []

    # 1) Replace pictures already inserted (matched by Alt Text / descr)
    # 2) Optionally replace text placeholders whose text equals an existing filename
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if text_mapping:
                replaced_text += _replace_text_in_shape(shape, text_mapping)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt = _get_shape_alt_text(shape)
                if alt:
                    candidate = images_dir / alt
                    if candidate.exists():
                        _replace_picture_image_in_place(slide, shape, candidate)
                        replaced_pictures += 1
                        replaced_files.append(alt)
                    else:
                        # If alt looks like an image filename, warn later.
                        if Path(alt).suffix.lower() in {".png", ".jpg", ".jpeg"}:
                            missing_files.append(alt)
                continue

            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            if allow_placeholder_text:
                candidate = images_dir / text
                if candidate.exists():
                    _replace_shape_with_picture(slide, shape, candidate)
                    replaced_placeholders += 1
                    replaced_files.append(text)
                else:
                    if Path(text).suffix.lower() in {".png", ".jpg", ".jpeg"}:
                        missing_files.append(text)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    # If input == output, write to a temp file first, then replace.
    if pptx_path.resolve() == output_path.resolve():
        tmp_path = output_path.with_suffix(output_path.suffix + ".tmp")
        prs.save(str(tmp_path))
        tmp_path.replace(output_path)
    else:
        prs.save(str(output_path))

    applied_text_keys = sorted(text_mapping.keys()) if text_mapping else []

    return (
        replaced_pictures,
        replaced_placeholders,
        replaced_text,
        replaced_files,
        missing_files,
        applied_text_keys,
    )


def _collect_pictures_alt_text(pptx_path: Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    alts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            alt = _get_shape_alt_text(shape)
            if alt:
                alts.append(alt)
    return alts


def _collect_text_placeholders(pptx_path: Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            t = (shape.text_frame.text or "").strip()
            if t:
                texts.append(t)
    return texts


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Atualiza um PPTX substituindo imagens existentes pelo Alt Text (descr) \
apontando para arquivos no diretório de imagens.\n\n"
            "Regra principal: se a imagem no PPT tiver Alt Text = '01_lucro_trimestres.png', \
e existir um arquivo com esse nome no diretório de imagens, o script troca a imagem \
mantendo posição/tamanho/crop."
        )
    )
    parser.add_argument(
        "--pptx",
        "--input",
        required=True,
        help="Caminho do PPTX de entrada (template que será atualizado).",
    )
    parser.add_argument(
        "--out",
        "--output",
        default=None,
        help=(
            "Caminho do PPTX de saída. Se omitido, cria '<entrada>.updated.pptx' ao lado do arquivo de entrada."
        ),
    )
    parser.add_argument(
        "--images-dir",
        default=".",
        help="Diretório onde estão os PNG/JPG a inserir (default: diretório atual).",
    )
    parser.add_argument(
        "--in-place",
        action="store_true",
        help="Sobrescreve o arquivo de entrada (cuidado).",
    )
    parser.add_argument(
        "--allow-placeholder-text",
        action="store_true",
        help=(
            "Também substitui caixas de texto cujo texto seja exatamente um nome de arquivo existente (ex: '02_lucro_9m.png')."
        ),
    )
    parser.add_argument(
        "--text-json",
        default=None,
        help=(
            "JSON com textos do LLM para preencher no PPT. Regra: substituir tokens '{{chave}}' ou shapes com Alt Text == chave. "
            "Aceita formato direto {titles:{...},subtitles:{...}} ou wrapper {response:{...}}."
        ),
    )
    return parser.parse_args()


def main() -> None:
    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
    args = _parse_args()

    pptx_path = Path(args.pptx).expanduser().resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX não encontrado: {pptx_path}")

    images_dir = Path(args.images_dir).expanduser().resolve()
    if not images_dir.exists():
        raise FileNotFoundError(f"Diretório de imagens não encontrado: {images_dir}")

    if args.in_place:
        output_path = pptx_path
    else:
        output_path = (
            Path(args.out).expanduser().resolve()
            if args.out
            else pptx_path.with_name(pptx_path.stem + ".updated" + pptx_path.suffix)
        )

    logging.info("PPTX: %s", pptx_path)
    logging.info("Imagens: %s", images_dir)
    logging.info("Saída: %s", output_path)

    (
        replaced_pictures,
        replaced_placeholders,
        replaced_text,
        replaced_files,
        missing_files,
        applied_text_keys,
    ) = update_presentation(
        pptx_path=pptx_path,
        output_path=output_path,
        images_dir=images_dir,
        allow_placeholder_text=bool(args.allow_placeholder_text),
        text_json=Path(args.text_json).expanduser().resolve() if args.text_json else None,
    )

    logging.info(
        "Substituições: pictures=%d placeholders=%d",
        replaced_pictures,
        replaced_placeholders,
    )

    if args.text_json:
        logging.info("Substituições: text=%d", replaced_text)
        if applied_text_keys:
            logging.info("Chaves de texto disponíveis (%d): %s", len(applied_text_keys), applied_text_keys)

    if replaced_files:
        logging.info("Arquivos aplicados (%d): %s", len(replaced_files), sorted(set(replaced_files)))

    # Verification summary
    alts = _collect_pictures_alt_text(output_path)
    print(f"OK: gerado {output_path}")
    print(f"VERIF: pictures={len(alts)} alts={sorted(set(alts))}")

    if args.text_json:
        texts = _collect_text_placeholders(output_path)
        remaining_tokens = [t for t in texts if "{{" in t and "}}" in t]
        print(f"VERIF: text_shapes={len(texts)} remaining_tokens={len(remaining_tokens)}")
        if remaining_tokens:
            print("VERIF: remaining_tokens_list=", sorted(set(remaining_tokens)))

    missing_unique = sorted(set(missing_files))
    if missing_unique:
        warnings.warn(
            "Faltam arquivos no diretório de imagens para alguns Alt Texts / placeholders: "
            + ", ".join(missing_unique),
            stacklevel=2,
        )


if __name__ == "__main__":
    main()
