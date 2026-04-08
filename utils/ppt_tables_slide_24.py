from __future__ import annotations

import argparse
from dataclasses import dataclass
from decimal import Decimal, ROUND_HALF_UP
from pathlib import Path

from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_boundaries
from pptx import Presentation


SLIDE24_TABLE_ALT_TEXT = "TABLE_SLIDE24_DRE"
SLIDE24_SHEET_NAME = "DRE Saida"
SLIDE24_HEADER_RANGE = "C3:K4"
SLIDE24_VALUES_RANGE = "C5:K18"
SLIDE24_FIXED_COORDS = frozenset({"C3"})
SLIDE24_VALUES_TARGET_START = (2, 0)


@dataclass(frozen=True)
class Slide24TableApplyResult:
    found: bool
    slide_index: int | None
    shape_name: str | None
    written_cells: int
    skipped_fixed_cells: int
    skipped_spanned_cells: int


def _cell_is_percent_formatted(cell) -> bool:
    try:
        fmt = cell.number_format
    except Exception:
        return False
    if not fmt:
        return False
    return "%" in str(fmt)


def _get_shape_alt_text(shape) -> str | None:
    try:
        cnv = shape._element.xpath(".//p:cNvPr")
        if cnv:
            return cnv[0].get("descr")
    except Exception:
        return None
    return None


def _copy_run_font(src_run, dst_run) -> None:
    if src_run is None:
        return
    try:
        src_font = src_run.font
        dst_font = dst_run.font
        dst_font.name = src_font.name
        dst_font.size = src_font.size
        dst_font.bold = src_font.bold
        dst_font.italic = src_font.italic
        dst_font.underline = src_font.underline
        try:
            dst_font.color.rgb = src_font.color.rgb
        except Exception:
            pass
    except Exception:
        return


def _set_cell_text_preserving_style(cell, text: str) -> None:
    text_frame = cell.text_frame
    first_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
    first_run = first_paragraph.runs[0] if first_paragraph and first_paragraph.runs else None
    paragraph_alignment = first_paragraph.alignment if first_paragraph is not None else None
    paragraph_level = first_paragraph.level if first_paragraph is not None else None

    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    if paragraph_alignment is not None:
        paragraph.alignment = paragraph_alignment
    if paragraph_level is not None:
        paragraph.level = paragraph_level
    run = paragraph.add_run()
    run.text = text
    _copy_run_font(first_run, run)


def _round_half_up(value: float, decimals: int) -> Decimal:
    quant = Decimal("1") if decimals <= 0 else Decimal("0." + ("0" * (decimals - 1)) + "1")
    return Decimal(str(float(value))).quantize(quant, rounding=ROUND_HALF_UP)


def _format_integer_accounting(value: float) -> str:
    rounded = int(_round_half_up(value, 0))
    rendered = f"{abs(rounded):,}".replace(",", ".")
    return f"({rendered})" if rounded < 0 else rendered


def _format_decimal_ptbr(value: float, decimals: int = 1) -> str:
    rounded = _round_half_up(value, decimals)
    return f"{rounded:.{decimals}f}".replace(".", ",")


def _coerce_text(value: object) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _format_slide24_cell(cell) -> str:
    value = cell.value
    if value is None:
        return ""

    if isinstance(value, bool):
        return str(value)

    row_idx = int(cell.row)
    col_idx = int(cell.column)

    if row_idx >= 5 and 4 <= col_idx <= 8 and isinstance(value, (int, float)):
        return _format_integer_accounting(float(value))

    if row_idx >= 5 and 9 <= col_idx <= 11 and isinstance(value, (int, float)):
        numeric_value = float(value)
        if _cell_is_percent_formatted(cell):
            numeric_value *= 100.0
        return _format_decimal_ptbr(numeric_value, 1)

    return _coerce_text(value)


def extract_slide24_table_block(
    *,
    xlsx_path: Path,
    sheet_name: str = SLIDE24_SHEET_NAME,
    source_range: str,
) -> list[list[str]]:
    wb = load_workbook(filename=xlsx_path, data_only=True)
    try:
        if sheet_name not in wb.sheetnames:
            raise ValueError(f"Aba não encontrada: {sheet_name!r}. Disponíveis: {wb.sheetnames}")
        ws = wb[sheet_name]
        min_col, min_row, max_col, max_row = range_boundaries(source_range)
        values: list[list[str]] = []
        for row in range(min_row, max_row + 1):
            rendered_row: list[str] = []
            for col in range(min_col, max_col + 1):
                rendered_row.append(_format_slide24_cell(ws.cell(row=row, column=col)))
            values.append(rendered_row)
        return values
    finally:
        wb.close()


def extract_slide24_table_headers(
    *,
    xlsx_path: Path,
    sheet_name: str = SLIDE24_SHEET_NAME,
    source_range: str = SLIDE24_HEADER_RANGE,
) -> list[list[str]]:
    return extract_slide24_table_block(
        xlsx_path=xlsx_path,
        sheet_name=sheet_name,
        source_range=source_range,
    )


def extract_slide24_table_values(
    *,
    xlsx_path: Path,
    sheet_name: str = SLIDE24_SHEET_NAME,
    source_range: str = SLIDE24_VALUES_RANGE,
) -> list[list[str]]:
    return extract_slide24_table_block(
        xlsx_path=xlsx_path,
        sheet_name=sheet_name,
        source_range=source_range,
    )


def _find_table_shape(prs: Presentation, *, table_alt_text: str):
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            if _get_shape_alt_text(shape) == table_alt_text:
                return slide_idx, shape
    return None


def _collect_writable_target_row_cells(table, row_idx: int):
    cells = []
    for col_idx in range(len(table.columns)):
        cell = table.cell(row_idx, col_idx)
        if getattr(cell, "is_spanned", False):
            continue
        cells.append(cell)
    return cells


def _apply_header_rows(
    *,
    table,
    values: list[list[str]],
    source_range: str,
) -> tuple[int, int, int]:
    min_col, min_row, _max_col, _max_row = range_boundaries(source_range)
    if len(table.rows) < len(values):
        raise ValueError(
            "Tabela do PowerPoint menor que o cabeçalho esperado do slide 24: "
            f"ppt_rows={len(table.rows)} header_rows={len(values)}"
        )

    written_cells = 0
    skipped_fixed_cells = 0
    skipped_spanned_cells = 0

    for row_offset, row_values in enumerate(values):
        source_items: list[str] = []
        for col_offset, rendered_value in enumerate(row_values):
            source_coord = f"{get_column_letter(min_col + col_offset)}{min_row + row_offset}"
            if source_coord in SLIDE24_FIXED_COORDS:
                skipped_fixed_cells += 1
                continue
            if rendered_value == "":
                continue
            source_items.append(rendered_value)

        target_cells = _collect_writable_target_row_cells(table, row_offset)
        row_spanned = len(table.columns) - len(target_cells)
        skipped_spanned_cells += max(0, row_spanned)

        if len(target_cells) < len(source_items):
            raise ValueError(
                "Cabeçalho do PowerPoint tem menos células visíveis do que o esperado: "
                f"row={row_offset} ppt_visible={len(target_cells)} source_values={len(source_items)}"
            )

        if len(target_cells) > len(source_items):
            # Em templates corporativos, o excedente tende a ficar à esquerda
            # por causa de células fixas/mescladas como o bloco do título.
            target_cells = target_cells[-len(source_items):] if source_items else []

        for target_cell, rendered_value in zip(target_cells, source_items):
            _set_cell_text_preserving_style(target_cell, rendered_value)
            written_cells += 1

    return written_cells, skipped_fixed_cells, skipped_spanned_cells


def _apply_table_block(
    *,
    table,
    values: list[list[str]],
    source_range: str,
    target_start_row: int,
    target_start_col: int,
    skip_empty: bool,
) -> tuple[int, int, int]:
    min_col, min_row, _max_col, _max_row = range_boundaries(source_range)
    expected_rows = target_start_row + len(values)
    expected_cols = target_start_col + (len(values[0]) if values else 0)
    if len(table.rows) < expected_rows or len(table.columns) < expected_cols:
        raise ValueError(
            "Tabela do PowerPoint menor que o bloco esperado do slide 24: "
            f"ppt={len(table.rows)}x{len(table.columns)} bloco={expected_rows}x{expected_cols}"
        )

    written_cells = 0
    skipped_fixed_cells = 0
    skipped_spanned_cells = 0

    for row_offset, row_values in enumerate(values):
        for col_offset, rendered_value in enumerate(row_values):
            source_coord = f"{get_column_letter(min_col + col_offset)}{min_row + row_offset}"
            if source_coord in SLIDE24_FIXED_COORDS:
                skipped_fixed_cells += 1
                continue
            if skip_empty and rendered_value == "":
                continue

            target_cell = table.cell(target_start_row + row_offset, target_start_col + col_offset)
            if getattr(target_cell, "is_spanned", False):
                skipped_spanned_cells += 1
                continue

            _set_cell_text_preserving_style(target_cell, rendered_value)
            written_cells += 1

    return written_cells, skipped_fixed_cells, skipped_spanned_cells


def apply_slide24_table_to_presentation(
    prs: Presentation,
    *,
    xlsx_path: Path,
    table_alt_text: str = SLIDE24_TABLE_ALT_TEXT,
    sheet_name: str = SLIDE24_SHEET_NAME,
    header_range: str = SLIDE24_HEADER_RANGE,
    values_range: str = SLIDE24_VALUES_RANGE,
    strict: bool = False,
) -> Slide24TableApplyResult:
    table_location = _find_table_shape(prs, table_alt_text=table_alt_text)
    if table_location is None:
        if strict:
            raise ValueError(f"Tabela com alt text {table_alt_text!r} não encontrada no PPT.")
        return Slide24TableApplyResult(
            found=False,
            slide_index=None,
            shape_name=None,
            written_cells=0,
            skipped_fixed_cells=0,
            skipped_spanned_cells=0,
        )

    header_values = extract_slide24_table_headers(
        xlsx_path=xlsx_path,
        sheet_name=sheet_name,
        source_range=header_range,
    )
    body_values = extract_slide24_table_values(
        xlsx_path=xlsx_path,
        sheet_name=sheet_name,
        source_range=values_range,
    )
    if not header_values and not body_values:
        raise ValueError("Ranges da tabela do slide 24 não retornaram valores.")

    slide_index, shape = table_location
    table = shape.table
    expected_rows = max(
        SLIDE24_VALUES_TARGET_START[0] + len(body_values),
        len(header_values),
    )
    expected_cols = max(
        SLIDE24_VALUES_TARGET_START[1] + (len(body_values[0]) if body_values else 0),
        len(header_values[0]) if header_values else 0,
    )
    if len(table.rows) < expected_rows or len(table.columns) < expected_cols:
        raise ValueError(
            "Tabela do PowerPoint menor que o range esperado do slide 24: "
            f"ppt={len(table.rows)}x{len(table.columns)} range={expected_rows}x{expected_cols}"
        )

    header_written, header_fixed, header_spanned = _apply_header_rows(
        table=table,
        values=header_values,
        source_range=header_range,
    )
    body_written, body_fixed, body_spanned = _apply_table_block(
        table=table,
        values=body_values,
        source_range=values_range,
        target_start_row=SLIDE24_VALUES_TARGET_START[0],
        target_start_col=SLIDE24_VALUES_TARGET_START[1],
        skip_empty=False,
    )

    return Slide24TableApplyResult(
        found=True,
        slide_index=slide_index,
        shape_name=getattr(shape, "name", None),
        written_cells=header_written + body_written,
        skipped_fixed_cells=header_fixed + body_fixed,
        skipped_spanned_cells=header_spanned + body_spanned,
    )


def apply_slide24_table_file(
    *,
    pptx_path: Path,
    output_path: Path,
    xlsx_path: Path,
    table_alt_text: str = SLIDE24_TABLE_ALT_TEXT,
    sheet_name: str = SLIDE24_SHEET_NAME,
    header_range: str = SLIDE24_HEADER_RANGE,
    values_range: str = SLIDE24_VALUES_RANGE,
    strict: bool = False,
) -> Slide24TableApplyResult:
    prs = Presentation(str(pptx_path))
    result = apply_slide24_table_to_presentation(
        prs,
        xlsx_path=xlsx_path,
        table_alt_text=table_alt_text,
        sheet_name=sheet_name,
        header_range=header_range,
        values_range=values_range,
        strict=strict,
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    if pptx_path.resolve() == output_path.resolve():
        tmp_path = output_path.with_suffix(output_path.suffix + ".tmp")
        prs.save(str(tmp_path))
        tmp_path.replace(output_path)
    else:
        prs.save(str(output_path))
    return result


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Preenche a tabela do slide 24 no PowerPoint a partir do workbook DRE Saida."
    )
    parser.add_argument("--pptx", required=True, help="PPTX de entrada contendo a tabela formatada.")
    parser.add_argument("--xlsx", required=True, help="Workbook de entrada com a aba DRE Saida.")
    parser.add_argument(
        "--output",
        help="PPTX de saída. Se omitido, grava '<entrada>.updated.pptx'.",
    )
    parser.add_argument(
        "--alt-text",
        default=SLIDE24_TABLE_ALT_TEXT,
        help=f"Alt text da tabela no PowerPoint. Default: {SLIDE24_TABLE_ALT_TEXT}",
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Falha se a tabela não for encontrada no PPT.",
    )
    return parser.parse_args()


def _main() -> None:
    args = _parse_args()
    pptx_path = Path(args.pptx).expanduser().resolve()
    xlsx_path = Path(args.xlsx).expanduser().resolve()
    output_path = (
        Path(args.output).expanduser().resolve()
        if args.output
        else pptx_path.with_name(pptx_path.stem + ".updated" + pptx_path.suffix)
    )

    result = apply_slide24_table_file(
        pptx_path=pptx_path,
        output_path=output_path,
        xlsx_path=xlsx_path,
        table_alt_text=str(args.alt_text),
        strict=bool(args.strict),
    )
    print(
        "slide24_table:",
        {
            "found": result.found,
            "slide_index": result.slide_index,
            "shape_name": result.shape_name,
            "written_cells": result.written_cells,
            "skipped_fixed_cells": result.skipped_fixed_cells,
            "skipped_spanned_cells": result.skipped_spanned_cells,
            "output_path": str(output_path),
        },
    )


if __name__ == "__main__":
    _main()
